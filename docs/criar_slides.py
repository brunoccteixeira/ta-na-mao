#!/usr/bin/env python3
"""
Script para criar/completar a apresentação PowerPoint do Tá na Mão
"""

import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_THEME_COLOR
except ImportError:
    print("Instalando python-pptx...")
    os.system("pip3 install python-pptx pillow")
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        from pptx.enum.dml import MSO_THEME_COLOR
    except ImportError:
        print("Erro: Não foi possível instalar python-pptx")
        sys.exit(1)

def criar_apresentacao_completa():
    """Cria a apresentação completa com todos os 14 slides"""
    
    # Criar nova apresentação
    prs = Presentation()
    print("Criando nova apresentação...")
    
    # Layout padrão
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    title_content_layout = prs.slide_layouts[1]  # Title and content
    
    # Cores do tema (verde-azul gradiente)
    cor_primaria = RGBColor(0, 102, 102)  # Verde escuro
    cor_secundaria = RGBColor(0, 153, 153)  # Verde água
    cor_destaque = RGBColor(255, 102, 0)  # Laranja para destaques
    
    # SLIDE 1: CAPA
    slide1 = prs.slides.add_slide(title_slide_layout)
    
    title = slide1.shapes.title
    title.text = "TÁ NA MÃO"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = cor_primaria
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle = slide1.placeholders[1]
    subtitle.text = "Última Milha de Direitos Públicos\n\nProposta de Piloto via CPSI / Espaço TEIA\n\nDe direito previsto em lei a direito na mão.\n\nGOVERNO FEDERAL | CAIXA | TÁ NA MÃO\n\nVersão 11.0 | Dezembro 2025"
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)
    
    # SLIDE 2: O PROBLEMA
    slide2 = prs.slides.add_slide(title_content_layout)
    slide2.shapes.title.text = "O PROBLEMA"
    
    content = slide2.placeholders[1].text_frame
    content.text = "O gargalo não é a verba. É a execução na ponta."
    
    p = content.add_paragraph()
    p.text = "\nO Brasil criou os programas. Garantiu os recursos. Mas o direito não chega."
    p.font.size = Pt(16)
    
    p = content.add_paragraph()
    p.text = "\n🍔 Analogia:"
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "Imagine se o iFood existisse, os restaurantes estivessem abertos, a comida pronta... mas a última milha fosse complexa demais para boa parte das famílias."
    p.font.size = Pt(14)
    p.font.italic = True
    
    p = content.add_paragraph()
    p.text = "\nÉ isso que acontece com R$ 50 bi em benefícios públicos."
    p.font.size = Pt(14)
    p.font.italic = True
    
    # SLIDE 3: VAZAMENTO SOCIAL  
    slide3 = prs.slides.add_slide(title_content_layout)
    slide3.shapes.title.text = "VAZAMENTO SOCIAL"
    
    content = slide3.placeholders[1].text_frame
    content.text = "Quem Tem Direito, Não Acessa"
    
    p = content.add_paragraph()
    p.text = "\n🔸 Até ~90%"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = cor_destaque
    
    p = content.add_paragraph()
    p.text = "Em programas específicos, a não-adesão chega a ~90% do público estimado"
    p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "\n💊 Farmácia Popular: Milhões com doenças crônicas têm direito. Penetração muito abaixo do potencial."
    p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "\n🩸 Dignidade Menstrual: ~24 mi elegíveis; >2 mi beneficiadas/ano (≈8%)"
    p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "\n⚡ TSEE: 7,7 mi famílias têm direito, mas não recebem. R$ 540/ano perdidos."
    p.font.size = Pt(14)
    
    # SLIDE 4: R$ 50 BI
    slide4 = prs.slides.add_slide(title_content_layout)
    slide4.shapes.title.text = "R$ 50 BI ESPERANDO CHEGAR NA PONTA"
    
    content = slide4.placeholders[1].text_frame
    content.text = "Do Saque Extraordinário ao Tá na Mão"
    
    p = content.add_paragraph()
    p.text = "\nO Saque Extraordinário do FGTS liberou R$ 23,6 bi para 32,7 mi trabalhadores."
    p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "\nProgramas abaixo já existem e somam ~R$ 50 bi esperando chegar na ponta:"
    p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "\n• PIS/PASEP (1971-1988): ~R$ 26,0 bi"
    p.font.size = Pt(16)
    
    p = content.add_paragraph()
    p.text = "• SVR/BCB (ainda não sacados): ~R$ 8,7 bi"
    p.font.size = Pt(16)
    
    p = content.add_paragraph()
    p.text = "• TSEE + Farmácia + Outros: ~R$ 15,3 bi"
    p.font.size = Pt(16)
    
    p = content.add_paragraph()
    p.text = "\n💰 TOTAL ESTIMADO: ~R$ 50 bi"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = cor_destaque
    
    # SLIDE 5: A SOLUÇÃO
    slide5 = prs.slides.add_slide(title_content_layout)
    slide5.shapes.title.text = "A SOLUÇÃO"
    
    content = slide5.placeholders[1].text_frame
    content.text = "Tá na Mão: Execução de Direitos via WhatsApp"
    
    p = content.add_paragraph()
    p.text = "\nInformação não basta. Executamos a jornada para que o direito seja ativado com rastreabilidade."
    p.font.size = Pt(16)
    
    p = content.add_paragraph()
    p.text = "\n📱 Canais atuais vs. Tá na Mão:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "\n• Autoatendimento → Execução assistida"
    p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "• Cidadão navega → Direito é ativado"
    p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "• Informação disponível → Jornada executada"
    p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "\n⭐ Regra de Ouro: Cidadão nunca paga pelo benefício público."
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = cor_destaque
    
    # SLIDE 6: MOTOR NORMATIVO
    slide6 = prs.slides.add_slide(title_content_layout)
    slide6.shapes.title.text = "MOTOR NORMATIVO"
    
    content = slide6.placeholders[1].text_frame
    content.text = "Inteligência de Elegibilidade e Execução"
    
    p = content.add_paragraph()
    p.text = "\n🔽 Motor de Elegibilidade"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "Cruza perfil do cidadão/família com catálogo vivo de benefícios"
    p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "\n🔗 Orquestração Multi-Sistemas"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "Jornada com integrações em CadÚnico, Gov.br e sistemas setoriais"
    p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "\n📋 Carteira de Direitos"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "Visão 360° persistente por família, não por transação isolada"
    p.font.size = Pt(14)
    
    # SLIDE 7: JORNADA DO CIDADÃO
    slide7 = prs.slides.add_slide(title_content_layout)
    slide7.shapes.title.text = "JORNADA DO CIDADÃO"
    
    content = slide7.placeholders[1].text_frame
    content.text = "Do Pedido à Ativação - 6 etapas"
    
    etapas = [
        "📱 Pedido: WhatsApp (CPF + CEP)",
        "✅ Confirmação: Consentimento LGPD",
        "🔍 Preparo: Motor cruza perfil com catálogo",
        "🛵 Execução: Orquestração com integrações oficiais",
        "📩 Ativado!: Comprovante do direito ativado",
        "💳 Histórico: Carteira atualizada"
    ]
    
    for i, etapa in enumerate(etapas, 1):
        p = content.add_paragraph()
        p.text = f"\n{i}. {etapa}"
        p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "\n🤝 Human-in-the-loop: validação humana em etapas críticas"
    p.font.size = Pt(12)
    p.font.italic = True
    
    # SLIDE 8: POR QUE A CAIXA
    slide8 = prs.slides.add_slide(title_content_layout)
    slide8.shapes.title.text = "POR QUE A CAIXA?"
    
    content = slide8.placeholders[1].text_frame
    content.text = "Parceiro Natural para Infraestrutura Nacional de Direitos"
    
    p = content.add_paragraph()
    p.text = "\n🏛️ Hub Histórico"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "Braço operacional de políticas sociais há décadas"
    p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "\n🤝 Confiança Institucional"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "Grau de confiança que nenhum banco privado possui"
    p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "\n🎯 Complementaridade"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "Não concorremos. Reforçamos o papel público da Caixa."
    p.font.size = Pt(14)
    
    # SLIDE 9: VALOR PARA A CAIXA
    slide9 = prs.slides.add_slide(title_content_layout)
    slide9.shapes.title.text = "VALOR PARA A CAIXA"
    
    content = slide9.placeholders[1].text_frame
    content.text = "O Que a Caixa Ganha"
    
    beneficios = [
        "📉 Menos Filas: Cidadão orientado via WhatsApp = menos atendimento presencial",
        "😊 Menos Reclamações: Direito ativado = menos frustração",
        "📱 Caixa Tem Fortalecido: WhatsApp complementa o app",
        "🏆 Protagonismo: Referência nacional em 'última milha' de direitos"
    ]
    
    for beneficio in beneficios:
        p = content.add_paragraph()
        p.text = f"\n• {beneficio}"
        p.font.size = Pt(14)
    
    # SLIDE 10: PILOTO
    slide10 = prs.slides.add_slide(title_content_layout)
    slide10.shapes.title.text = "PILOTO"
    
    content = slide10.placeholders[1].text_frame
    content.text = "Validação em Dois Cenários"
    
    p = content.add_paragraph()
    p.text = "\n🏙️ Cenário 1: Teste de Carga (Volume)"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "Região Metropolitana de São Paulo - validar robustez técnica"
    p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "\n🌾 Cenário 2: Teste de Acesso (Vulnerabilidade)"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "Município estratégico no Nordeste - validar impacto social"
    p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "\n📊 Escopo: 6-9 meses | ~60k famílias | TSEE + Farmácia + PIS/PASEP"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = cor_destaque
    
    # SLIDE 11: KPIs
    slide11 = prs.slides.add_slide(title_content_layout)
    slide11.shapes.title.text = "KPIs"
    
    content = slide11.placeholders[1].text_frame
    content.text = "Laboratório de Eficiência e Impacto"
    
    kpis = [
        "💰 Custo de Ativação: R$ por família (WhatsApp vs. agência)",
        "📈 Valor Acessado: R$ em direitos efetivamente acessados",
        "✅ Taxa de Conclusão: % elegíveis que completam jornada",
        "📉 Redução de Filas: Δ atendimentos nas agências piloto"
    ]
    
    for kpi in kpis:
        p = content.add_paragraph()
        p.text = f"\n• {kpi}"
        p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "\n🎯 Potencial: R$ 2.000-4.000/ano por família elegível"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = cor_destaque
    
    # SLIDE 12: O QUE PRECISAMOS DA CAIXA
    slide12 = prs.slides.add_slide(title_content_layout)
    slide12.shapes.title.text = "O QUE PRECISAMOS DA CAIXA"
    
    content = slide12.placeholders[1].text_frame
    content.text = "Parceria no Piloto"
    
    necessidades = [
        "👥 Squad Conjunto: TEIA + Benefícios + Caixa Tem + Compliance",
        "📊 Dados Agregados: Indicadores de impacto (sem dados individuais)",
        "📣 Co-comunicação: Canal oficial verificado + Caixa Tem",
        "🔒 Segurança & Compliance: LGPD + auditoria + anti-golpes"
    ]
    
    for necessidade in necessidades:
        p = content.add_paragraph()
        p.text = f"\n• {necessidade}"
        p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "\n⚠️ NUNCA pedimos senha Gov.br | NUNCA pedimos PIX"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 0, 0)
    
    # SLIDE 13: CPSI/TEIA
    slide13 = prs.slides.add_slide(title_content_layout)
    slide13.shapes.title.text = "ENQUADRAMENTO: CPSI / ESPAÇO TEIA"
    
    content = slide13.placeholders[1].text_frame
    content.text = "Modelo de Inovação que a Caixa Já Utiliza"
    
    p = content.add_paragraph()
    p.text = "\n📋 Marco Legal: LC nº 182/2021 (Marco Legal das Startups)"
    p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "\n💰 Envelope: até R$ 1,6 milhão"
    p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "\n⏱️ Vigência: até 12 meses, prorrogável uma vez"
    p.font.size = Pt(14)
    
    p = content.add_paragraph()
    p.text = "\n🚀 Gatilho de Escala: Em caso de sucesso, contrato ampliado (3-5x)"
    p.font.size = Pt(14)
    p.font.color.rgb = cor_destaque
    
    p = content.add_paragraph()
    p.text = "\n🤝 Governança: Comitê conjunto com dashboards compartilhados"
    p.font.size = Pt(14)
    
    # SLIDE 14: FECHAMENTO/CTA
    slide14 = prs.slides.add_slide(title_content_layout)
    slide14.shapes.title.text = "TRANSFORMAR DIREITO EM REALIDADE"
    
    content = slide14.placeholders[1].text_frame
    content.text = "Com a Caixa no centro"
    
    p = content.add_paragraph()
    p.text = "\nO Brasil já criou os programas e garantiu os recursos."
    p.font.size = Pt(16)
    
    p = content.add_paragraph()
    p.text = "\nO Tá na Mão propõe a infraestrutura de última milha que transforma burocracia em clique."
    p.font.size = Pt(16)
    
    p = content.add_paragraph()
    p.text = "\n💰 ~R$ 50 bi esperando chegar na ponta."
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = cor_destaque
    p.alignment = PP_ALIGN.CENTER
    
    p = content.add_paragraph()
    p.text = "\n🤝 Queremos construir isso junto com a Caixa."
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "\n🎯 PRÓXIMO PASSO:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "Workshop de 2h com TEIA, Benefícios e Compliance"
    p.font.size = Pt(16)
    
    # Salvar apresentação
    filename = "Ta_na_Mao_Apresentacao_Caixa_v11_COMPLETA.pptx"
    prs.save(filename)
    print(f"\n✅ Apresentação completa salva como: {filename}")
    print("📊 14 slides criados com sucesso!")
    
    return filename

if __name__ == "__main__":
    criar_apresentacao_completa()