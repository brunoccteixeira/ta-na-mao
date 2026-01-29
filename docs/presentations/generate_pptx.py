#!/usr/bin/env python3
"""
Gerador de Apresentacao Ta na Mao - CAIXA v2
Estilo: Dark theme (#1a1a2e) + Verde limao (#c8ff00)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# Cores do tema
DARK_BG = RGBColor(0x1a, 0x1a, 0x2e)
DARK_BG_LIGHT = RGBColor(0x25, 0x27, 0x3c)
VERDE_LIMAO = RGBColor(0xc8, 0xff, 0x00)
BRANCO = RGBColor(0xff, 0xff, 0xff)
CINZA_CLARO = RGBColor(0xaa, 0xaa, 0xaa)
AZUL_CARD = RGBColor(0x2a, 0x2d, 0x45)

# Dimensoes
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_background(slide, color):
    """Define cor de fundo do slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, top=Inches(0.5), left=Inches(0.5), width=Inches(12),
              font_size=44, color=VERDE_LIMAO, bold=True):
    """Adiciona titulo ao slide"""
    shape = slide.shapes.add_textbox(left, top, width, Inches(1))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Arial"
    return shape


def add_text(slide, text, top, left, width, height=Inches(1), font_size=18,
             color=BRANCO, bold=False, align=PP_ALIGN.LEFT):
    """Adiciona texto ao slide"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Arial"
    p.alignment = align
    return shape


def add_bullet_text(slide, items, top, left, width, font_size=16, color=BRANCO):
    """Adiciona lista com bullets"""
    shape = slide.shapes.add_textbox(left, top, width, Inches(3))
    tf = shape.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = Pt(8)
    return shape


def add_card(slide, title, content, left, top, width=Inches(2.8), height=Inches(2),
             title_color=VERDE_LIMAO, bg_color=AZUL_CARD):
    """Adiciona card com titulo e conteudo"""
    # Retangulo de fundo
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    # Titulo do card
    title_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15),
                                          width - Inches(0.3), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = "Arial"

    # Conteudo do card
    content_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.55),
                                            width - Inches(0.3), height - Inches(0.7))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(12)
    p.font.color.rgb = BRANCO
    p.font.name = "Arial"


def add_stat_card(slide, number, label, left, top, width=Inches(2.5), height=Inches(1.5)):
    """Adiciona card de estatistica"""
    # Retangulo verde limao
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = VERDE_LIMAO
    shape.line.fill.background()

    # Numero grande
    num_box = slide.shapes.add_textbox(left, top + Inches(0.2), width, Inches(0.8))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    p.font.bold = True
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

    # Label
    label_box = slide.shapes.add_textbox(left, top + Inches(0.9), width, Inches(0.5))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER


def add_table_simple(slide, headers, rows, left, top, col_widths):
    """Adiciona tabela simples"""
    num_rows = len(rows) + 1
    num_cols = len(headers)

    table = slide.shapes.add_table(num_rows, num_cols, left, top,
                                    sum(col_widths), Inches(0.4 * num_rows)).table

    # Ajustar larguras das colunas
    for i, width in enumerate(col_widths):
        table.columns[i].width = width

    # Header
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = VERDE_LIMAO
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = DARK_BG
        p.font.name = "Arial"

    # Rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = AZUL_CARD
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = BRANCO
            p.font.name = "Arial"


def create_presentation():
    """Cria a apresentacao completa"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # Layout em branco

    # ==================== SLIDE 1: CAPA ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BG)

    # Titulo principal
    add_title(slide, "TA NA MAO", top=Inches(2.5), left=Inches(0.5),
              width=Inches(12), font_size=72, color=VERDE_LIMAO)

    # Subtitulo
    add_text(slide, "Ultimo Km de Direitos Publicos", top=Inches(3.5),
             left=Inches(0.5), width=Inches(12), font_size=32, color=BRANCO, bold=True)

    # Tagline
    add_text(slide, "De direito previsto em lei a direito na mao.", top=Inches(4.3),
             left=Inches(0.5), width=Inches(12), font_size=20, color=VERDE_LIMAO)

    # Rodape
    add_text(slide, "GOVERNO FEDERAL | CAIXA | TA NA MAO", top=Inches(6.3),
             left=Inches(0.5), width=Inches(12), font_size=12, color=CINZA_CLARO,
             align=PP_ALIGN.CENTER)
    add_text(slide, "PROPOSTA DE PILOTO VIA CPSI / ESPACO TEIA | JANEIRO 2026",
             top=Inches(6.6), left=Inches(0.5), width=Inches(12), font_size=10,
             color=CINZA_CLARO, align=PP_ALIGN.CENTER)

    # ==================== SLIDE 2: O GARGALO ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BG)

    add_title(slide, "O gargalo nao e a verba. E a execucao na ponta.")

    add_text(slide, "O Brasil criou os programas. Garantiu os recursos.\nMas o direito nem sempre chega.",
             top=Inches(1.5), left=Inches(0.5), width=Inches(6), font_size=20, color=BRANCO)

    # Quote box
    quote_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(0.5), Inches(2.5), Inches(6), Inches(1.2))
    quote_shape.fill.solid()
    quote_shape.fill.fore_color.rgb = AZUL_CARD
    quote_shape.line.fill.background()

    add_text(slide, '"Imagine se o iFood existisse, a comida estivesse pronta... mas o ultimo km fosse complexo demais. E isso que acontece com R$ 50 bi em beneficios."',
             top=Inches(2.6), left=Inches(0.7), width=Inches(5.6), font_size=14,
             color=CINZA_CLARO)

    # Bullets
    items = [
        "Cidadao como despachante digital: Navega dezenas de apps e senhas",
        "Decifra regras federais, estaduais, municipais ao mesmo tempo",
        "O que falta nao e politica publica - e infraestrutura de execucao"
    ]
    add_bullet_text(slide, items, top=Inches(4), left=Inches(0.5), width=Inches(6))

    # ==================== SLIDE 3: VAZAMENTO SOCIAL ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BG)

    add_title(slide, "Vazamento Social: Quem Tem Direito, Nao Acessa")

    # Grande numero
    add_stat_card(slide, "~90%", "Em programas especificos,\na nao-adesao chega a 90%\ndo publico estimado.",
                  left=Inches(0.5), top=Inches(1.5), width=Inches(3.5), height=Inches(2.5))

    # Cards dos programas
    add_card(slide, "Tarifa Social (TSEE)",
             "7,7 milhoes de familias tem direito, mas nao recebem.\nDeixam de economizar R$ 540/ano.",
             left=Inches(4.5), top=Inches(1.5), width=Inches(4), height=Inches(1.3))

    add_card(slide, "Farmacia Popular",
             "Milhoes com doencas cronicas tem direito, mas a penetracao efetiva e muito abaixo do potencial.",
             left=Inches(4.5), top=Inches(3), width=Inches(4), height=Inches(1.3))

    add_card(slide, "Dignidade Menstrual",
             "Publico estimado ~24 mi; ~2,5 mi beneficiadas/ano (apenas 10% acessam).",
             left=Inches(4.5), top=Inches(4.5), width=Inches(4), height=Inches(1.3))

    add_text(slide, "Fontes: ANEEL 2025; MDS 2025; MS 2025", top=Inches(6.2),
             left=Inches(0.5), width=Inches(8), font_size=10, color=CINZA_CLARO)

    # ==================== SLIDE 4: OPORTUNIDADE ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BG)

    add_title(slide, "Do Saque Extraordinario ao Ta na Mao")

    add_text(slide, "O Saque Extraordinario 2022 liberou R$ 23,6 bi. Os programas abaixo ja existem, somam ~R$ 50 bi e precisam de infraestrutura de execucao, nao nova MP.",
             top=Inches(1.4), left=Inches(0.5), width=Inches(12), font_size=16, color=BRANCO)

    # Barras de valor
    # PIS/PASEP
    bar1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(3.5), Inches(2.5), Inches(7), Inches(0.6))
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = VERDE_LIMAO
    bar1.line.fill.background()
    add_text(slide, "PIS/PASEP (Cotas 1971-1988)", top=Inches(2.5), left=Inches(0.5),
             width=Inches(3), font_size=14, color=BRANCO)
    add_text(slide, "~R$ 26,3 Bilhoes", top=Inches(2.55), left=Inches(3.7),
             width=Inches(3), font_size=14, color=DARK_BG, bold=True)

    # TSEE + Farmacia
    bar2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(3.5), Inches(3.4), Inches(4.5), Inches(0.6))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = VERDE_LIMAO
    bar2.line.fill.background()
    add_text(slide, "TSEE + Farmacia Popular + Outros", top=Inches(3.4), left=Inches(0.5),
             width=Inches(3), font_size=14, color=BRANCO)
    add_text(slide, "~R$ 15 Bilhoes", top=Inches(3.45), left=Inches(3.7),
             width=Inches(3), font_size=14, color=DARK_BG, bold=True)

    # SVR
    bar3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(3.5), Inches(4.3), Inches(3), Inches(0.6))
    bar3.fill.solid()
    bar3.fill.fore_color.rgb = VERDE_LIMAO
    bar3.line.fill.background()
    add_text(slide, "Valores a Receber (SVR)", top=Inches(4.3), left=Inches(0.5),
             width=Inches(3), font_size=14, color=BRANCO)
    add_text(slide, "~R$ 8-10 Bilhoes", top=Inches(4.35), left=Inches(3.7),
             width=Inches(3), font_size=14, color=DARK_BG, bold=True)

    add_text(slide, "Transformar recursos ja autorizados em renda e dignidade na ponta.",
             top=Inches(5.5), left=Inches(0.5), width=Inches(12), font_size=18,
             color=VERDE_LIMAO, bold=True, align=PP_ALIGN.CENTER)

    # ==================== SLIDE 5: A SOLUCAO ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BG)

    add_title(slide, "Ta na Mao: Execucao de Direitos via WhatsApp")

    add_text(slide, "Informacao nao basta. Executamos a jornada para que o direito seja ativado.",
             top=Inches(1.4), left=Inches(0.5), width=Inches(6), font_size=18, color=BRANCO)

    # Quote
    quote_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(0.5), Inches(2.2), Inches(6), Inches(1.3))
    quote_shape.fill.solid()
    quote_shape.fill.fore_color.rgb = AZUL_CARD
    quote_shape.line.fill.background()

    add_text(slide, '"A partir de poucos dados, descobrimos a quais beneficios o cidadao ja tem direito. Um Motor Normativo orquestra a burocracia ate que o direito esteja, de fato, na mao."',
             top=Inches(2.35), left=Inches(0.7), width=Inches(5.6), font_size=13, color=CINZA_CLARO)

    # Tabela comparativa
    add_table_simple(slide,
                     ["Canais Atuais (Gov.br)", "Ta na Mao"],
                     [
                         ["Autoatendimento digital", "Execucao assistida"],
                         ["Cidadao navega", "Direito e ativado"],
                         ["Informacao disponivel", "Jornada executada"],
                         ["Transacao isolada", "Visao 360 da familia"]
                     ],
                     left=Inches(7), top=Inches(1.8),
                     col_widths=[Inches(2.5), Inches(2.5)])

    add_text(slide, "Regra de Ouro: Cidadao nunca paga pelo beneficio.",
             top=Inches(5.5), left=Inches(0.5), width=Inches(12), font_size=16,
             color=VERDE_LIMAO, bold=True, align=PP_ALIGN.CENTER)

    # ==================== SLIDE 6: O QUE JA CONSTRUIMOS ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BG)

    add_title(slide, "O Que Ja Construimos")

    add_text(slide, "Produto Funcional - Pronto para Piloto",
             top=Inches(1.4), left=Inches(0.5), width=Inches(12), font_size=20,
             color=BRANCO, align=PP_ALIGN.CENTER)

    # 4 cards de metricas
    add_stat_card(slide, "16", "Telas no App Android",
                  left=Inches(0.8), top=Inches(2.2), width=Inches(2.7), height=Inches(1.8))
    add_stat_card(slide, "25", "Tools do Agente IA",
                  left=Inches(3.8), top=Inches(2.2), width=Inches(2.7), height=Inches(1.8))
    add_stat_card(slide, "10", "Programas Rastreados",
                  left=Inches(6.8), top=Inches(2.2), width=Inches(2.7), height=Inches(1.8))
    add_stat_card(slide, "100%", "Backend Async",
                  left=Inches(9.8), top=Inches(2.2), width=Inches(2.7), height=Inches(1.8))

    # Tecnologias
    add_text(slide, "Tecnologia:", top=Inches(4.5), left=Inches(0.5), width=Inches(12),
             font_size=16, color=VERDE_LIMAO, bold=True)

    items = [
        "Google Gemini 2.0 Flash (IA conversacional)",
        "Arquitetura Orchestrator + Sub-agents",
        "Integracao WhatsApp (Twilio)",
        "Dashboard web para gestores"
    ]
    add_bullet_text(slide, items, top=Inches(5), left=Inches(0.5), width=Inches(6), font_size=14)

    # ==================== SLIDE 7: TELAS DO APP ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BG)

    add_title(slide, "App Ta na Mao - Interface Real", font_size=36)

    # 4 mockups de telas (representados como cards)
    screens = [
        ("Home", "Resumo de beneficios\nDinheiro esquecido\nProximos pagamentos\nServicos perto de voce"),
        ("Chat com IA", "Assistente inteligente\nQuick actions\nEnvio de fotos\nRespostas estruturadas"),
        ("Carteira", "Beneficios ativos\nElegiveis\nHistorico completo"),
        ("Mapa", "CRAS proximos\nFarmacias populares\nIntegracao GPS")
    ]

    for i, (title, content) in enumerate(screens):
        left = Inches(0.5 + i * 3.1)

        # Frame do celular
        phone_frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              left, Inches(1.5), Inches(2.8), Inches(5))
        phone_frame.fill.solid()
        phone_frame.fill.fore_color.rgb = RGBColor(0x35, 0x38, 0x50)
        phone_frame.line.color.rgb = VERDE_LIMAO
        phone_frame.line.width = Pt(2)

        # Titulo da tela
        add_text(slide, title, top=Inches(1.7), left=left + Inches(0.1),
                 width=Inches(2.6), font_size=16, color=VERDE_LIMAO, bold=True,
                 align=PP_ALIGN.CENTER)

        # Conteudo
        content_box = slide.shapes.add_textbox(left + Inches(0.15), Inches(2.3),
                                                Inches(2.5), Inches(3.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for line in content.split('\n'):
            p = tf.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(11)
            p.font.color.rgb = BRANCO
            p.font.name = "Arial"
            p.space_after = Pt(6)

    # ==================== SLIDE 8: MOTOR NORMATIVO ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BG)

    add_title(slide, "Motor Normativo: Inteligencia de Elegibilidade")

    add_text(slide, "Interpretacao de regras + orquestracao de jornadas.",
             top=Inches(1.4), left=Inches(0.5), width=Inches(12), font_size=16,
             color=CINZA_CLARO, align=PP_ALIGN.CENTER)

    # 4 cards
    cards = [
        ("Motor de Elegibilidade", "Cruza perfil com catalogo vivo de beneficios (Uniao, Estados, Municipios). Nao e atendimento, e execucao assistida."),
        ("Orquestracao", "Integracoes e fluxos oficiais em CadUnico e Gov.br, com suporte humano em etapas criticas."),
        ("Carteira de Direitos", "Visao 360 persistente por familia, historico completo de beneficios ativados."),
        ("Atualizacao Continua", "Quando as regras mudam (lei, decreto), o motor e atualizado. Catalogo vivo.")
    ]

    for i, (title, content) in enumerate(cards):
        row = i // 2
        col = i % 2
        left = Inches(0.5 + col * 6.2)
        top = Inches(2 + row * 2.3)
        add_card(slide, title, content, left=left, top=top, width=Inches(5.8), height=Inches(2))

    # ==================== SLIDE 9: POR QUE A CAIXA ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BG)

    add_title(slide, "Por Que a Caixa e o Parceiro Natural?")

    add_text(slide, "A juncao da capilaridade da Caixa com o WhatsApp cria uma Infraestrutura Nacional de Execucao.",
             top=Inches(1.5), left=Inches(0.5), width=Inches(6), font_size=18, color=BRANCO)

    items = [
        "Hub Historico: A Caixa e o braco operacional de politicas sociais ha decadas.",
        "Confianca: Grau de confianca do cidadao que nenhum banco privado possui.",
        "Complementaridade: Cada cidadao atendido via WhatsApp e uma pessoa a menos na fila da agencia."
    ]
    add_bullet_text(slide, items, top=Inches(2.8), left=Inches(0.5), width=Inches(6), font_size=16)

    # ==================== SLIDE 10: O QUE A CAIXA GANHA ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BG)

    add_title(slide, "O Que a Caixa Ganha")

    cards = [
        ("Menos Filas", "Cidadao orientado e com direito ativado via WhatsApp = menos atendimento presencial."),
        ("Menos Reclamacoes", "Direito ativado = menos frustracao = menos Ouvidoria e SAC."),
        ("Caixa Tem Fortalecido", "WhatsApp complementa o app, ampliando alcance para quem tem dificuldade digital."),
        ("Protagonismo Social", "Caixa como referencia nacional em \"ultima milha\" de direitos.")
    ]

    for i, (title, content) in enumerate(cards):
        row = i // 2
        col = i % 2
        left = Inches(0.5 + col * 6.2)
        top = Inches(1.8 + row * 2.5)
        add_card(slide, title, content, left=left, top=top, width=Inches(5.8), height=Inches(2.2))

    # ==================== SLIDE 11: PILOTO PROPOSTO ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BG)

    add_title(slide, "Piloto: Validacao em Sao Paulo")

    add_text(slide, "Escopo: 6-9 meses | ~1.000 familias | Farmacia Popular",
             top=Inches(1.4), left=Inches(0.5), width=Inches(12), font_size=16,
             color=CINZA_CLARO, align=PP_ALIGN.CENTER)

    # Card Beneficio Hero
    add_card(slide, "Beneficio-Heroi: Farmacia Popular",
             "• Alto impacto individual (R$ 100-150/mes)\n• Milhoes com doencas cronicas\n• Jornada de ativacao bem definida\n• Validacao rapida de eficiencia",
             left=Inches(0.5), top=Inches(2), width=Inches(5.5), height=Inches(2.5))

    # KPIs
    add_text(slide, "KPIs do Piloto:", top=Inches(2), left=Inches(6.5),
             width=Inches(6), font_size=16, color=VERDE_LIMAO, bold=True)

    add_table_simple(slide,
                     ["KPI", "Meta"],
                     [
                         ["Custo de Ativacao", "< 50% do custo atual"],
                         ["Valor Acessado", "> R$ 800/familia/ano"],
                         ["Taxa de Conclusao", "> 70%"],
                         ["NPS", "> 50"]
                     ],
                     left=Inches(6.5), top=Inches(2.5),
                     col_widths=[Inches(2.5), Inches(2.5)])

    add_text(slide, "Enquadramento: CPSI (LC 182/2021) | Ate R$ 1,6M | Via Espaco TEIA",
             top=Inches(5.5), left=Inches(0.5), width=Inches(12), font_size=14,
             color=CINZA_CLARO, align=PP_ALIGN.CENTER)

    # ==================== SLIDE 12: FECHAMENTO ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BG)

    add_text(slide, "Transformar direito em realidade", top=Inches(2),
             left=Inches(0.5), width=Inches(12), font_size=44, color=VERDE_LIMAO,
             bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, "- com a Caixa no centro.", top=Inches(2.8),
             left=Inches(0.5), width=Inches(12), font_size=44, color=VERDE_LIMAO,
             bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, 'O Brasil ja criou os programas e garantiu os recursos.\nO "Ta na Mao" propoe a infraestrutura do ultimo km\nque transforma burocracia em clique.',
             top=Inches(4), left=Inches(0.5), width=Inches(12), font_size=18,
             color=BRANCO, align=PP_ALIGN.CENTER)

    add_text(slide, "~R$ 50 bi esperando chegar na ponta.",
             top=Inches(5), left=Inches(0.5), width=Inches(12), font_size=16,
             color=VERDE_LIMAO, align=PP_ALIGN.CENTER)

    # CTA Box
    cta_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(3.5), Inches(5.8), Inches(6.3), Inches(0.8))
    cta_shape.fill.solid()
    cta_shape.fill.fore_color.rgb = AZUL_CARD
    cta_shape.line.color.rgb = VERDE_LIMAO
    cta_shape.line.width = Pt(1)

    add_text(slide, "Proximo Passo: Workshop de 2h com TEIA, Beneficios e Compliance.",
             top=Inches(5.95), left=Inches(3.5), width=Inches(6.3), font_size=14,
             color=BRANCO, bold=True, align=PP_ALIGN.CENTER)

    return prs


if __name__ == "__main__":
    print("Gerando apresentacao Ta na Mao - CAIXA v2...")
    prs = create_presentation()

    output_path = os.path.join(os.path.dirname(__file__), "Ta_na_Mao_CAIXA_v2.pptx")
    prs.save(output_path)

    print(f"Apresentacao salva em: {output_path}")
    print("12 slides criados com sucesso!")
